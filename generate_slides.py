from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

# Colors
C_ISABELLINE = RGBColor(248, 242, 237) # #F8F2ED
C_COBALT = RGBColor(48, 80, 162)       # #3050A2
C_MEADOW = RGBColor(41, 177, 152)      # #29B198
C_WHITE = RGBColor(255, 255, 255)
C_DARK = RGBColor(25, 32, 39)          # #192027

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_ISABELLINE

def add_title(slide, text, color=C_COBALT):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.name = 'Inter'
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_glass_card(slide, left, top, width, height, text="", font_size=16, bg_color=C_WHITE, font_color=C_DARK, bold_title=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(230, 230, 230)
    shape.line.width = Pt(1)
    
    # Adjust corner radius to be soft
    shape.adjustments[0] = 0.15 
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Inter'
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold_title
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow(slide, left, top, width, height, rotation=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_MEADOW
    shape.line.fill.background()
    shape.rotation = rotation

prs = Presentation()
# Set 16:9 ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# ----------------- SLIDE 1 -----------------
s1 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s1)
# Main glass card
c1 = add_glass_card(s1, 1.5, 1.5, 10.333, 4.5)
# Title
t1 = s1.shapes.add_textbox(Inches(2), Inches(2), Inches(9.333), Inches(1))
p1 = t1.text_frame.add_paragraph()
p1.text = "HỆ THỐNG TỰ ĐỘNG HÓA LÊN GIÁO ÁN\n& ĐIỀU PHỐI HỌC CỤ"
p1.font.bold = True
p1.font.size = Pt(36)
p1.font.name = 'Inter'
p1.font.color.rgb = C_COBALT
p1.alignment = PP_ALIGN.CENTER

t2 = s1.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(1.5))
t2.text_frame.word_wrap = True
p2 = t2.text_frame.add_paragraph()
p2.text = "Trích xuất tài nguyên học tập (như Twinkl) từ Cloud, tự động thiết kế giáo án tích hợp kịch bản hoạt động góc và soạn bản tin tóm tắt trên Drive+Excel gửi phụ huynh+giáo viên qua Zalo. Tự động bóc tách danh mục vật liệu thủ công, role-play và nhắc nhở chuẩn bị trước 1 tuần."
p2.font.size = Pt(16)
p2.font.name = 'Inter'
p2.font.color.rgb = C_DARK
p2.alignment = PP_ALIGN.CENTER

t3 = s1.shapes.add_textbox(Inches(2), Inches(5.2), Inches(9.333), Inches(0.5))
p3 = t3.text_frame.add_paragraph()
p3.text = "Học viên: Hồng Nhung  |  Nhóm 1"
p3.font.bold = True
p3.font.size = Pt(18)
p3.font.color.rgb = C_MEADOW
p3.alignment = PP_ALIGN.CENTER

# ----------------- SLIDE 2 -----------------
s2 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s2)
add_title(s2, "BỐI CẢNH & QUY TRÌNH NGHIỆP VỤ")
steps = ["1. Học liệu thô\n(PDF, Twinkl)", "2. Phân bổ\nLịch trình", "3. Biên soạn\nGiáo án V3", "4. Kiểm duyệt\nChất lượng", "5. Triển khai\n& Nhắc nhở"]
for i, step in enumerate(steps):
    add_glass_card(s2, 0.8 + i*2.4, 3, 2, 2, step, font_size=16, bold_title=True)
    if i < 4:
        add_arrow(s2, 2.9 + i*2.4, 3.8, 0.2, 0.3)
add_glass_card(s2, 1, 5.5, 11.333, 1, "Quy trình Academic Quality Control: Cốt lõi duy trì tiêu chuẩn sư phạm mầm non", 18, C_MEADOW, C_WHITE, True)

# ----------------- SLIDE 3 -----------------
s3 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s3)
add_title(s3, "THÁCH THỨC & MỤC TIÊU (KPIs)")
add_glass_card(s3, 1, 2, 4.5, 4.5, bg_color=RGBColor(255, 245, 245))
t = s3.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(4), Inches(4))
tf = t.text_frame
tf.word_wrap = True
tf.text = "THÁCH THỨC HIỆN TẠI\n\n• Mất 4-5 tiếng/tuần/lớp để soạn bài.\n• Dễ nhầm lẫn độ khó giữa trẻ 3-4 và 5-6 tuổi.\n• Quên ghi chú an toàn gây rủi ro cao."
for p in tf.paragraphs:
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 50, 50)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(22)

add_glass_card(s3, 6.5, 2, 6, 4.5, bg_color=C_WHITE)
t2 = s3.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(4))
tf2 = t2.text_frame
tf2.word_wrap = True
tf2.text = "MỤC TIÊU SCOPE (KPIs)\n\n• TỐC ĐỘ: Giảm thời gian từ 4-5 tiếng xuống < 3 PHÚT.\n• CHẤT LƯỢNG: Chuẩn hóa 100% Montessori, đủ Part 1-4.\n• NHẬN DIỆN: Đúng tone màu và font chữ Inter."
for p in tf2.paragraphs:
    p.font.size = Pt(18)
    p.font.color.rgb = C_COBALT
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.size = Pt(22)

# ----------------- SLIDE 4 -----------------
s4 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s4)
add_title(s4, "PHÂN CẤP HỆ THỐNG: ĐIỀU PHỐI – QUẢN LÝ – THỰC THI")
add_glass_card(s4, 3, 2, 7.333, 1.2, "1. ĐIỀU PHỐI (Master Pipeline)\nTổng chỉ huy, tiếp nhận lệnh và điều hướng luồng AI", 16, C_COBALT, C_WHITE, True)
add_arrow(s4, 6.5, 3.3, 0.3, 0.4, rotation=90)
add_glass_card(s4, 2, 3.8, 9.333, 1.2, "2. QUẢN LÝ (Planners & Classifiers)\nPhân tích độ tuổi, phân loại môn học & xếp lịch 2 tuần thông minh", 16, C_MEADOW, C_WHITE, True)
add_arrow(s4, 6.5, 5.1, 0.3, 0.4, rotation=90)
add_glass_card(s4, 1, 5.6, 11.333, 1.2, "3. THỰC THI (Writers & Sync Tools)\nĐọc PDF thô, sinh giáo án LLM, Render DOCX chuẩn màu & Sync Google Drive", 16, RGBColor(220, 235, 250), C_DARK, True)

# ----------------- SLIDE 5 -----------------
s5 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s5)
add_title(s5, "KIẾN TRÚC 7 THÀNH TỐ AGENTIC WORKSPACE")
add_glass_card(s5, 5.666, 3, 2, 2, "AGENTIC\nWORKSPACE", 20, C_COBALT, C_WHITE, True)
positions = [
    (2, 1.5, "1. Agent Core (MICRO)"), (9.3, 1.5, "2. Workflow (OIPO)"),
    (1, 3.2, "3. Rules (CLEAR)"), (10.3, 3.2, "4. Knowledge Base (RAG)"),
    (2, 5.5, "5. Google Sync"), (9.3, 5.5, "6. Human Checkpoint"),
    (5.666, 6, "7. Handoff (JSON/MD)")
]
for p in positions:
    add_glass_card(s5, p[0], p[1], 2, 1, p[2], 12, C_WHITE, C_DARK, True)

# ----------------- SLIDE 6 -----------------
s6 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s6)
add_title(s6, "QUY TRÌNH VẬN HÀNH & LUỒNG DỮ LIỆU")
add_glass_card(s6, 1, 2.5, 11.333, 4, bg_color=C_WHITE)
steps_flow = ["Data Converter\n(Trích xuất MD)", "Curriculum Planner\n(Lập lịch 10 tiết)", "Lesson Plan Writer\n(Viết bài 4 Parts)", "DOCX Renderer\n(Định dạng File)", "Workspace Sync\n(Up Drive/Sheet)"]
for i, step in enumerate(steps_flow):
    add_glass_card(s6, 1.3 + i*2.2, 3.5, 1.9, 1.5, step, 14, C_MEADOW if i%2==0 else C_COBALT, C_WHITE, True)
    if i < 4:
        add_arrow(s6, 3.3 + i*2.2, 4.1, 0.1, 0.2)
t = s6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1))
t.text_frame.text = "Lệnh kích hoạt tự động hóa: Đọc. [Độ tuổi]"
t.text_frame.paragraphs[0].font.size = Pt(18)
t.text_frame.paragraphs[0].font.color.rgb = C_DARK
t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ----------------- SLIDE 7 -----------------
s7 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s7)
add_title(s7, "ĐIỂM CHẠM CON NGƯỜI (HUMAN CHECKPOINT)")
add_glass_card(s7, 1.5, 2.5, 4, 3, "AI Tự Động Hóa\n\n- Soạn giáo án xong\n- Chờ duyệt để upload", 16)
add_glass_card(s7, 7.8, 2.5, 4, 3, "Google Cloud\n\n- Lưu Drive\n- Báo cáo Sheet", 16)
add_glass_card(s7, 5.8, 2.8, 1.7, 2.4, "Quality Gate\n(Manager)", 16, C_MEADOW, C_WHITE, True)
add_glass_card(s7, 3, 6, 7.333, 1, "Gõ 'Duyệt.' -> Chuyển tiếp || Hoặc Góp ý -> Lưu feedback.txt cho AI sửa lại", 14, RGBColor(220, 240, 220), C_DARK)

# ----------------- SLIDE 8 -----------------
s8 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s8)
add_title(s8, "HỆ THỐNG ĐÁNH GIÁ TỰ ĐỘNG (AI GRADING)")
t = s8.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(0.8))
t.text_frame.text = "Barem chấm điểm 100 điểm với 7 Tiêu chí Khắt khe"
t.text_frame.paragraphs[0].font.size = Pt(20)
t.text_frame.paragraphs[0].font.bold = True
t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_glass_card(s8, 1.5, 2.8, 3, 3, "1. Xử lý dữ liệu (15)\n2. Chuẩn Montessori (20)\n3. Theo Độ Tuổi (15)\n4. Scaffolding (15)", 14)
add_glass_card(s8, 5, 2.8, 3, 3, "5. Độ an toàn (10)\n6. Hứng thú TPR (15)\n7. Định dạng (10)", 14)
add_glass_card(s8, 8.5, 2.8, 3.5, 3, "⚠ RED FLAGS & ALERT\n\nNếu <50đ hoặc thiếu an toàn: Gửi Mail Cảnh báo Tự động ngay lập tức.", 14, RGBColor(255, 230, 230), RGBColor(200, 0, 0), True)

# ----------------- SLIDE 9 -----------------
s9 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s9)
add_title(s9, "ACADEMIC MANAGER DASHBOARD")
add_glass_card(s9, 1.5, 2, 10.333, 4.5, bg_color=RGBColor(240, 245, 255))
t = s9.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(3.5))
tf = t.text_frame
tf.word_wrap = True
tf.text = "Giao diện Glassmorphism trên Streamlit (Port 9009)\n\n• Biểu diễn trực quan điểm số dưới dạng Progress bar.\n• Hiển thị chi tiết lỗi (Red flags) & Đề xuất sửa đổi tự động.\n• Quản lý tập trung toàn bộ tiến độ làm giáo án 2 tuần."
for p in tf.paragraphs:
    p.font.size = Pt(20)
    p.font.color.rgb = C_COBALT
    p.space_after = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(24)

# ----------------- SLIDE 10 -----------------
s10 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s10)
add_title(s10, "KẾT QUẢ & LIVE DEMO THỰC TẾ")
add_glass_card(s10, 1.5, 2.5, 3, 2, "100%\nTuân thủ An toàn\n& Montessori", 20, C_MEADOW, C_WHITE, True)
add_glass_card(s10, 5.166, 2.5, 3, 2, "< 3 PHÚT\nHoàn thành\n10 giáo án/2 tuần", 20, C_COBALT, C_WHITE, True)
add_glass_card(s10, 8.833, 2.5, 3, 2, "0 GIỜ\nĐiều phối viên\nphải làm tay", 20, RGBColor(255, 150, 50), C_WHITE, True)
add_glass_card(s10, 3.5, 5.5, 6.333, 1, "BẮT ĐẦU LIVE DEMO", 24, C_DARK, C_WHITE, True)

prs.save("Presentation_Nhom1.pptx")
print("Presentation generated successfully at Presentation_Nhom1.pptx")
